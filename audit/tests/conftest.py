"""
conftest.py — Фикстуры и генерация тестовых данных для всех тестов верификации.
"""
import os
import sys
import pytest
from pathlib import Path

# Добавляем scripts в PYTHONPATH
SCRIPTS_DIR = Path(__file__).resolve().parents[2] / ".claude" / "skills" / "verification" / "scripts"
sys.path.insert(0, str(SCRIPTS_DIR))

TEST_DATA_DIR = Path(__file__).parent / "test_data"
TEST_DATA_DIR.mkdir(exist_ok=True)


@pytest.fixture(scope="session", autouse=True)
def generate_test_files():
    """Генерирует все тестовые файлы один раз перед запуском тестов."""
    _generate_all()
    yield


def _generate_all():
    """Программная генерация всех тестовых файлов."""
    _gen_good_docx()
    _gen_bad_docx()
    _gen_empty_docx()
    _gen_good_xlsx()
    _gen_bad_xlsx()
    _gen_empty_xlsx()
    _gen_docx_with_refs_ok()
    _gen_docx_with_refs_broken()
    _gen_docx_with_numbering_ok()
    _gen_docx_with_numbering_bad()
    _gen_docx_with_dates()
    _gen_good_pptx()
    _gen_bad_pptx()
    _gen_html_matching_pptx()
    _gen_html_mismatched()
    _gen_cross_file_docx()
    _gen_cross_file_xlsx()
    _gen_docx_v1()
    _gen_docx_v2()


def _gen_good_docx():
    from docx import Document
    from docx.shared import Pt, Cm, RGBColor
    doc = Document()
    for section in doc.sections:
        section.top_margin = Cm(2.0)
        section.bottom_margin = Cm(2.0)
        section.left_margin = Cm(3.0)
        section.right_margin = Cm(1.5)
    # Заголовок
    h = doc.add_heading("Тестовый документ", level=1)
    for run in h.runs:
        run.font.size = Pt(22)
        run.font.color.rgb = RGBColor(0x00, 0x70, 0xC0)
    # Обычный абзац
    p = doc.add_paragraph()
    run = p.add_run("Это тестовый абзац с правильным форматированием.")
    run.font.name = "Times New Roman"
    run.font.size = Pt(14)
    p.style = doc.styles["Normal"]
    doc.save(str(TEST_DATA_DIR / "good.docx"))


def _gen_bad_docx():
    from docx import Document
    from docx.shared import Pt, Cm
    doc = Document()
    for section in doc.sections:
        section.top_margin = Cm(1.0)  # неверно
        section.bottom_margin = Cm(1.0)  # неверно
        section.left_margin = Cm(1.0)  # неверно
        section.right_margin = Cm(1.0)  # неверно
    p = doc.add_paragraph()
    run = p.add_run("Текст с неправильным шрифтом.")
    run.font.name = "Arial"
    run.font.size = Pt(10)
    p.style = doc.styles["Normal"]
    doc.save(str(TEST_DATA_DIR / "bad.docx"))


def _gen_empty_docx():
    from docx import Document
    doc = Document()
    doc.save(str(TEST_DATA_DIR / "empty.docx"))


def _gen_good_xlsx():
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Данные"
    ws["A1"] = "Наименование"
    ws["B1"] = "Сумма"
    ws["C1"] = "% доля"
    ws["A2"] = "Позиция 1"
    ws["B2"] = 100
    ws["C2"] = 25
    ws["A3"] = "Позиция 2"
    ws["B3"] = 200
    ws["C3"] = 50
    ws["A4"] = "Позиция 3"
    ws["B4"] = 100
    ws["C4"] = 25
    ws["A5"] = "Итого"
    ws["B5"] = 400  # верная сумма
    ws["C5"] = 100
    wb.save(str(TEST_DATA_DIR / "good.xlsx"))


def _gen_bad_xlsx():
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Данные"
    ws["A1"] = "Наименование"
    ws["B1"] = "Бюджет"
    ws["C1"] = "% доля"
    ws["A2"] = "Позиция 1"
    ws["B2"] = 100
    ws["C2"] = 25
    ws["A3"] = "Позиция 2"
    ws["B3"] = 200
    ws["C3"] = 50
    ws["A4"] = "Позиция 3"
    ws["B4"] = 100
    ws["C4"] = 25
    ws["A5"] = "Итого"
    ws["B5"] = 999  # неверная сумма!
    ws["C5"] = 150  # > 100%!
    ws["A6"] = "Спец"
    ws["B6"] = -50  # отрицательное значение в финансовом столбце
    wb.save(str(TEST_DATA_DIR / "bad.xlsx"))


def _gen_empty_xlsx():
    import openpyxl
    wb = openpyxl.Workbook()
    wb.save(str(TEST_DATA_DIR / "empty.xlsx"))


def _gen_docx_with_refs_ok():
    from docx import Document
    doc = Document()
    doc.add_paragraph("1. Введение")
    doc.add_paragraph("1.1 Область применения")
    doc.add_paragraph("Таблица 1")
    doc.add_paragraph("Рисунок 1")
    doc.add_paragraph("См. Таблица 1 и рисунок 1.")
    doc.add_paragraph("См. п. 1.1")
    # Добавляем таблицу
    doc.add_table(rows=2, cols=2)
    doc.save(str(TEST_DATA_DIR / "refs_ok.docx"))


def _gen_docx_with_refs_broken():
    from docx import Document
    doc = Document()
    doc.add_paragraph("1. Введение")
    doc.add_paragraph("Таблица 1")
    # Ссылка на несуществующие объекты
    doc.add_paragraph("См. Таблица 5 и рисунок 3.")
    doc.add_paragraph("См. п. 2.1")
    doc.add_paragraph("Приложение Б")
    doc.add_paragraph("См. приложение В")  # нет такого
    doc.add_table(rows=2, cols=2)  # только 1 таблица
    doc.save(str(TEST_DATA_DIR / "refs_broken.docx"))


def _gen_docx_with_numbering_ok():
    from docx import Document
    doc = Document()
    doc.add_paragraph("Таблица 1. Первая")
    doc.add_paragraph("Таблица 2. Вторая")
    doc.add_paragraph("Таблица 3. Третья")
    doc.add_paragraph("Рисунок 1. Первый")
    doc.add_paragraph("Рисунок 2. Второй")
    doc.save(str(TEST_DATA_DIR / "numbering_ok.docx"))


def _gen_docx_with_numbering_bad():
    from docx import Document
    doc = Document()
    doc.add_paragraph("Таблица 1. Первая")
    doc.add_paragraph("Таблица 1. Дубликат!")  # дубликат
    doc.add_paragraph("Таблица 3. Третья")  # пропуск 2
    doc.add_paragraph("Рисунок 1. Первый")
    doc.add_paragraph("Рисунок 3. Третий")  # пропуск 2
    doc.save(str(TEST_DATA_DIR / "numbering_bad.docx"))


def _gen_docx_with_dates():
    from docx import Document
    doc = Document()
    doc.add_paragraph("Документ от 15 марта 2025 года.")
    doc.add_paragraph("Дата подписания: 01.04.2025")
    doc.add_paragraph("ISO формат: 2025-06-15")
    doc.add_paragraph("Старая дата: 01.01.1980")  # подозрительная
    doc.add_paragraph("Далёкое будущее: 15.06.2035")  # подозрительная
    doc.save(str(TEST_DATA_DIR / "dates.docx"))


def _gen_good_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    prs = Presentation()
    # 16:9
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Тестовый слайд"
    # Добавляем ноту
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Это спикер-нота."

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    notes2 = slide2.notes_slide
    notes2.notes_text_frame.text = "Вторая нота."
    prs.save(str(TEST_DATA_DIR / "good.pptx"))


def _gen_bad_pptx():
    from pptx import Presentation
    from pptx.util import Emu, Pt
    prs = Presentation()
    # 4:3 - неверное соотношение
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(6858000)
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Плохой слайд"
    # Маленький шрифт
    from pptx.util import Inches
    from pptx.util import Pt as PPt
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Мелкий текст"
    run.font.size = PPt(8)
    # Без нот
    prs.save(str(TEST_DATA_DIR / "bad.pptx"))


def _gen_html_matching_pptx():
    html = """<!DOCTYPE html>
<html><body>
<section><h1>Тестовый слайд</h1><p>Содержимое первого слайда.</p></section>
<section><h1>Второй слайд</h1><p>Содержимое второго слайда.</p></section>
</body></html>"""
    (TEST_DATA_DIR / "matching.html").write_text(html, encoding="utf-8")


def _gen_html_mismatched():
    html = """<!DOCTYPE html>
<html><body>
<section><h1>Совсем другой контент</h1><p>Ничего общего с презентацией.</p></section>
</body></html>"""
    (TEST_DATA_DIR / "mismatched.html").write_text(html, encoding="utf-8")


def _gen_cross_file_docx():
    from docx import Document
    doc = Document()
    doc.add_paragraph("Бюджет проекта составляет 1 500 000 рублей.")
    doc.add_paragraph("Штат: 25 человек.")
    doc.add_paragraph("Руководитель: Иванов И.И.")
    doc.save(str(TEST_DATA_DIR / "cross.docx"))


def _gen_cross_file_xlsx():
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Штат"
    ws["A1"] = "Бюджет"
    ws["B1"] = "Штат чел."
    ws["A2"] = 1500000
    ws["B2"] = 30  # расхождение! в docx 25
    wb.save(str(TEST_DATA_DIR / "cross.xlsx"))


def _gen_docx_v1():
    from docx import Document
    doc = Document()
    doc.add_paragraph("Версия 1 документа.")
    doc.add_paragraph("Этот абзац будет изменён.")
    doc.add_paragraph("Этот абзац останется.")
    doc.save(str(TEST_DATA_DIR / "v1.docx"))


def _gen_docx_v2():
    from docx import Document
    doc = Document()
    doc.add_paragraph("Версия 2 документа.")
    doc.add_paragraph("Этот абзац был изменён!")
    doc.add_paragraph("Этот абзац останется.")
    doc.add_paragraph("Добавлен новый абзац.")
    doc.save(str(TEST_DATA_DIR / "v2.docx"))
