#!/usr/bin/env python3
"""
verify_pptx_format.py — Проверка формата презентации (16:9, шрифты, ноты).
"""

import argparse, json, sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt, Emu
except ImportError:
    print("pip install python-pptx --break-system-packages", file=sys.stderr)
    sys.exit(1)

ALLOWED_FONTS = {"Arial", "Calibri", "Helvetica", "Helvetica Neue", "Segoe UI", "Tahoma", "Verdana"}


def verify(filepath, allowed_fonts=None):
    prs = Presentation(filepath)
    findings = []
    fonts_whitelist = allowed_fonts or ALLOWED_FONTS

    # Проверка соотношения сторон
    w, h = prs.slide_width, prs.slide_height
    ratio = w / h if h else 0
    expected_ratio = 16 / 9  # 1.778
    if abs(ratio - expected_ratio) > 0.05:
        findings.append({
            "severity": "warning",
            "location": "Презентация",
            "expected": "16:9 (1.778)",
            "actual": f"{ratio:.3f}",
            "description": f"Соотношение сторон {ratio:.2f}, ожидалось 16:9",
        })

    # Проверка слайдов
    min_font = 100  # начальное значение для поиска минимума
    slides_without_notes = []
    non_standard_fonts = set()
    for i, slide in enumerate(prs.slides, 1):
        # Шрифты: размер + название
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            pt = run.font.size / Pt(1)  # EMU to pt
                            if pt < min_font and pt > 0:
                                min_font = pt
                        if run.font.name and run.font.name not in fonts_whitelist:
                            non_standard_fonts.add(run.font.name)
        
        # Спикер-ноты
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if not notes_text:
                slides_without_notes.append(i)
        else:
            slides_without_notes.append(i)

    if min_font < 12:
        findings.append({
            "severity": "warning",
            "location": "Шрифты",
            "expected": ">= 12pt",
            "actual": f"{min_font:.0f}pt",
            "description": f"Минимальный шрифт {min_font:.0f}pt — может быть нечитаемым при проекции",
        })

    if non_standard_fonts:
        findings.append({
            "severity": "info",
            "location": "Шрифты",
            "expected": f"из списка: {', '.join(sorted(fonts_whitelist))}",
            "actual": ', '.join(sorted(non_standard_fonts)),
            "description": f"Нестандартные шрифты: {', '.join(sorted(non_standard_fonts))}",
        })

    if slides_without_notes:
        if len(slides_without_notes) > len(prs.slides) * 0.5:
            findings.append({
                "severity": "info",
                "location": "Спикер-ноты",
                "expected": "на каждом слайде",
                "actual": f"отсутствуют на {len(slides_without_notes)} из {len(prs.slides)}",
                "description": f"Слайды без нот: {slides_without_notes[:10]}{'...' if len(slides_without_notes)>10 else ''}",
            })

    items_checked = len(prs.slides) + 2  # slides + ratio + fonts
    items_warned = len([f for f in findings if f["severity"] == "warning"])
    status = "warn" if items_warned > 0 else "pass"

    return {
        "script": "verify_pptx_format.py",
        "status": status,
        "details": f"Слайдов: {len(prs.slides)}, соотношение: {ratio:.2f}, мин. шрифт: {min_font:.0f}pt",
        "items_checked": items_checked,
        "items_passed": items_checked - items_warned,
        "items_warned": items_warned,
        "items_failed": 0,
        "findings": findings,
    }

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("files", nargs="+")
    parser.add_argument("--json", action="store_true")
    args = parser.parse_args()
    for f in args.files:
        if f.endswith(".pptx"):
            result = verify(f)
            print(json.dumps(result, ensure_ascii=False, indent=2) if args.json else f"📊 {f}: {result['status']}")
