#!/usr/bin/env python3
"""
verify_pptx_html_sync.py — Согласованность pptx и html версий презентации.
Сравнивает тексты слайдов pptx с секциями html.
"""

import argparse, json, re, sys
from pathlib import Path

def extract_pptx_text(filepath):
    from pptx import Presentation
    prs = Presentation(filepath)
    slides = []
    for slide in prs.slides:
        text_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_parts.append(shape.text_frame.text.strip())
        slides.append(" ".join(text_parts))
    return slides

def extract_html_text(filepath):
    text = Path(filepath).read_text(encoding="utf-8")
    # Убираем скрипты и стили
    text = re.sub(r'<script[^>]*>.*?</script>', '', text, flags=re.DOTALL)
    text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL)
    # Разделяем по section/slide-маркерам
    sections = re.split(r'<section|<div[^>]*class="[^"]*slide[^"]*"', text, flags=re.IGNORECASE)
    result = []
    for s in sections[1:]:  # пропускаем до первой секции
        clean = re.sub(r'<[^>]+>', ' ', s)
        clean = re.sub(r'\s+', ' ', clean).strip()
        if clean:
            result.append(clean[:500])  # ограничиваем длину
    return result

def normalize(text):
    return re.sub(r'\s+', ' ', text.lower().strip())

def verify(files):
    pptx_file = next((f for f in files if f.endswith(".pptx")), None)
    html_file = next((f for f in files if f.endswith(".html")), None)
    
    if not pptx_file or not html_file:
        return {"script": "verify_pptx_html_sync.py", "status": "skip",
                "details": "Нужны оба файла: .pptx и .html", "items_checked": 0,
                "items_passed": 0, "items_warned": 0, "items_failed": 0, "findings": []}

    pptx_slides = extract_pptx_text(pptx_file)
    html_sections = extract_html_text(html_file)
    findings = []

    # Проверка количества
    if len(pptx_slides) != len(html_sections):
        findings.append({
            "severity": "warning",
            "location": "Количество слайдов/секций",
            "expected": f"{len(pptx_slides)} (pptx)",
            "actual": f"{len(html_sections)} (html)",
            "description": f"Количество не совпадает: pptx={len(pptx_slides)}, html={len(html_sections)}",
        })

    # Попарное сравнение (по минимальному)
    for i in range(min(len(pptx_slides), len(html_sections))):
        pptx_norm = normalize(pptx_slides[i])
        html_norm = normalize(html_sections[i])
        
        # Извлекаем ключевые слова (>3 символов)
        pptx_words = set(w for w in pptx_norm.split() if len(w) > 3)
        html_words = set(w for w in html_norm.split() if len(w) > 3)
        
        if pptx_words and html_words:
            overlap = len(pptx_words & html_words) / max(len(pptx_words), len(html_words))
            if overlap < 0.3:
                findings.append({
                    "severity": "warning",
                    "location": f"Слайд/секция {i+1}",
                    "expected": "совпадение >30%",
                    "actual": f"{overlap:.0%} совпадение",
                    "description": f"Низкое совпадение текста между pptx и html на позиции {i+1}",
                })

    items_checked = min(len(pptx_slides), len(html_sections)) + 1
    items_warned = len(findings)
    status = "warn" if items_warned > 0 else "pass"

    return {
        "script": "verify_pptx_html_sync.py",
        "status": status,
        "details": f"pptx: {len(pptx_slides)} слайдов, html: {len(html_sections)} секций",
        "items_checked": items_checked,
        "items_passed": items_checked - items_warned,
        "items_warned": items_warned,
        "items_failed": 0,
        "findings": findings,
    }

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("files", nargs="+")
    parser.add_argument("--json", action="store_true")
    args = parser.parse_args()
    result = verify(args.files)
    print(json.dumps(result, ensure_ascii=False, indent=2) if args.json else f"🔄 Sync: {result['status']}")
